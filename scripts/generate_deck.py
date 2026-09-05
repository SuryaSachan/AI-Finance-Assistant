"""Generate clean, high-impact PowerPoint presentation with exact architecture diagram."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation(output_path="AI_Finance_Assistant_Demo.pptx"):
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme colors
    COLOR_BG = RGBColor(15, 23, 42)          # Slate 900
    COLOR_CARD = RGBColor(30, 41, 59)        # Slate 800
    COLOR_CARD_BORDER = RGBColor(51, 65, 85) # Slate 700
    COLOR_ACCENT = RGBColor(56, 189, 248)    # Sky 400
    COLOR_ACCENT2 = RGBColor(129, 140, 248)  # Indigo 400
    COLOR_SUCCESS = RGBColor(52, 211, 153)   # Emerald 400
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252)# Slate 50
    COLOR_TEXT_MUTED = RGBColor(203, 213, 225)# Slate 300
    COLOR_HIGHLIGHT = RGBColor(251, 191, 36) # Amber 400
    COLOR_AMBER = RGBColor(251, 191, 36)     # Amber 400
    COLOR_ROSE = RGBColor(244, 63, 94)       # Rose 500

    def add_bg(slide, color=COLOR_BG):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, tag, title, subtitle):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        p0.text = tag.upper()
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_ACCENT

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT_MAIN
        p1.space_before = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 1: PROBLEM
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    add_bg(s1)
    add_header(s1, "Slide 1 · The Challenge", "Why Traditional GenAI Fails in Financial Analytics",
               "Finance requires 100% precision. Standard LLMs are probabilistic, prone to hallucinations, and lack auditability.")

    problems = [
        ("01. The Hallucination Hazard",
         "LLMs invent plausible numbers, miscalculate sums, and fabricate transactions. In finance, an answer that is 'roughly right' is completely unacceptable.",
         COLOR_ACCENT,
         ["• LLMs cannot do reliable multi-step math",
          "• Unchecked outputs create compliance & balance sheet risks",
          "• Generative text conceals subtle data fabrication"]),
        ("02. The Security & Cost Barrier",
         "Sending enterprise ledger data to cloud frontier models (GPT-4) creates massive privacy exposure and unsustainable token bills.",
         COLOR_ACCENT2,
         ["• High latency and recurring token costs",
          "• Sensitive vendor/account data exposed to external APIs",
          "• Complete dependency on internet availability"]),
        ("03. The 'Black Box' Trust Deficit",
         "When an AI claims 'spend was ₹4.2 Cr', auditors and CFOs cannot verify how the number was derived or what filters were applied.",
         COLOR_HIGHLIGHT,
         ["• Zero lineage back to underlying records",
          "• No inspectable SQL or query plan",
          "• Inability to prove reconciliation or spot anomalies"])
    ]

    card_w = Inches(3.64)
    card_h = Inches(4.8)
    card_y = Inches(1.9)

    for i, (p_title, p_desc, p_accent, p_bullets) in enumerate(problems):
        card_x = Inches(0.8 + i * 4.0)
        shape = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_CARD_BORDER
        shape.line.width = Pt(1.5)

        tb = s1.shapes.add_textbox(card_x + Inches(0.25), card_y + Inches(0.25), card_w - Inches(0.5), card_h - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = p_accent

        p_desc_para = tf.add_paragraph()
        p_desc_para.text = p_desc
        p_desc_para.font.size = Pt(12)
        p_desc_para.font.color.rgb = COLOR_TEXT_MUTED
        p_desc_para.space_before = Pt(10)
        p_desc_para.space_after = Pt(14)

        for b in p_bullets:
            bp = tf.add_paragraph()
            bp.text = b
            bp.font.size = Pt(11)
            bp.font.color.rgb = COLOR_TEXT_MAIN
            bp.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 2: HIGH LEVEL DESIGN (EXACT ARCHITECTURE DIAGRAM - DARK MODE)
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_bg(s2, COLOR_BG) # Sleek dark slate background matching deck theme

    # Insert the high-resolution dark mode architecture diagram
    diag_path = "architecture_high_level_design.png"
    if os.path.exists(diag_path):
        s2.shapes.add_picture(
            diag_path,
            left=Inches(0.4),
            top=Inches(0.35),
            width=Inches(12.533),
            height=Inches(6.8)
        )

    # -------------------------------------------------------------
    # SLIDE 3: MODEL CHOICE RATIONALE
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_bg(s3)
    add_header(s3, "Slide 3 · Model Choice Rationale", "Selected Model: Qwen2.5-3B-Instruct — Lowest Size, Maximum Reliability",
               "Why a compact 3B local model outperforms bloated frontier models for enterprise financial inquiry.")

    left_x = Inches(0.8)
    left_y = Inches(1.9)
    left_w = Inches(6.8)
    left_h = Inches(4.9)

    left_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, left_y, left_w, left_h)
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLOR_CARD
    left_box.line.color.rgb = COLOR_CARD_BORDER
    left_box.line.width = Pt(1.5)

    tb_left = s3.shapes.add_textbox(left_x + Inches(0.3), left_y + Inches(0.25), left_w - Inches(0.6), left_h - Inches(0.5))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True

    p = tf_left.paragraphs[0]
    p.text = "Strategic Advantages of Qwen2.5-3B-Instruct"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    pillars = [
        ("1. Narrow Responsibility by Design",
         "The LLM is only tasked with slot-filling into a 10-key JSON schema and phrasing final text. It is NEVER asked to write arbitrary SQL or compute arithmetic."),
        ("2. Ultra-Lightweight Local Execution",
         "Consumes just ~2GB RAM at 4-bit quantisation. Runs effortlessly on CPU/laptop via Ollama with zero external cloud dependencies or API costs."),
        ("3. Superior JSON Adherence",
         "Outperformed comparable models (Llama-3.2-3B, Phi-3.5-mini) in adhering strictly to JSON formatting without verbose, unparseable conversational filler."),
        ("4. Tier 3 Deterministic Safety Net",
         "The architecture includes a rule-based engine that achieves 20/20 on accuracy benchmarks even when the LLM is completely offline. AI expands phrasing, not accuracy.")
    ]

    for title, desc in pillars:
        pt = tf_left.add_paragraph()
        pt.text = title
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN
        pt.space_before = Pt(8)

        pd = tf_left.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = COLOR_TEXT_MUTED
        pd.space_before = Pt(2)

    right_x = Inches(7.85)
    right_y = Inches(1.9)
    right_w = Inches(4.68)
    right_h = Inches(4.9)

    right_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, right_y, right_w, right_h)
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_CARD
    right_box.line.color.rgb = COLOR_CARD_BORDER
    right_box.line.width = Pt(1.5)

    tb_right = s3.shapes.add_textbox(right_x + Inches(0.3), right_y + Inches(0.25), right_w - Inches(0.6), right_h - Inches(0.5))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True

    pr = tf_right.paragraphs[0]
    pr.text = "Benchmark & Efficiency Metrics"
    pr.font.size = Pt(17)
    pr.font.bold = True
    pr.font.color.rgb = COLOR_SUCCESS

    metrics = [
        ("Parameter Size", "3 Billion (6× below the 20B cap)"),
        ("Benchmark Accuracy", "20 / 20 (100% on ground-truth SQL test)"),
        ("Prompt Footprint", "≈ 550 tokens (compact ~350 token schema)"),
        ("Completion Budget", "≈ 120 tokens (fast, deterministic JSON)"),
        ("Query Latency", "3 – 65 ms on DuckDB (over 250k records)"),
        ("Deployment Mode", "Local Ollama / CPU, zero cloud data leak"),
        ("Model Swappability", "Compatible with OpenAI/Sarvam drop-in")
    ]

    for label, val in metrics:
        p_row = tf_right.add_paragraph()
        p_row.text = f"• {label}: "
        p_row.font.size = Pt(11)
        p_row.font.bold = True
        p_row.font.color.rgb = COLOR_TEXT_MAIN
        p_row.space_before = Pt(6)
        
        run = p_row.add_run()
        run.text = val
        run.font.bold = False
        run.font.color.rgb = COLOR_ACCENT

    # -------------------------------------------------------------
    # SLIDE 4: RECENT ROBUSTNESS & SECURITY ENHANCEMENTS
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_bg(s4)
    add_header(s4, "Slide 4 · Latest Production Upgrades", "Recent Robustness, Guardrails & Performance Optimizations",
               "Summary of latest defense-in-depth changes added to ensure production resilience.")

    upgrades = [
        ("🛡️ Prompt Injection & Social Engineering Defense",
         "Added regex filters in planner.py to instantly reject simulated persona attacks ('pretend you have access to passwords') and raw SQL keywords (DROP TABLE, DELETE FROM).",
         COLOR_ROSE,
         ["• Rejection happens before LLM or DB execution",
          "• Clear explanatory refusal rather than unhandled crash",
          "• Prevents jailbreak attempts and schema leaks"]),
        ("📅 Out-of-Range Period Guardrail",
         "Pre-execution date span check in engine.py: verifies whether requested periods fall within known database boundaries before firing DB queries.",
         COLOR_AMBER,
         ["• Automatically tells the user exact available date range",
          "• Prevents confusing 'no data found' for out-of-bounds dates",
          "• Resolves anchor date dynamically from DuckDB"]),
        ("⚡ Monthly Rollup Pre-Aggregation",
         "Introduced v_rollup_monthly pre-aggregated view in DuckDB for high-level monthly trends and aggregations.",
         COLOR_SUCCESS,
         ["• Drastically reduces query scan size over massive datasets",
          "• Automatic routing in sql_builder.py for eligible queries",
          "• Keeps response latency consistently sub-50ms"]),
        ("👻 Ghost-Filter Elimination in Multi-Turn",
         "Prevents internal identifiers (e.g. UUID account_id) from leaking into subsequent unrelated questions.",
         COLOR_ACCENT,
         ["• Strips ghost filters unless question explicitly restates account",
          "• Handles period transitions cleanly across multi-turn chat",
          "• Retains 100% (20/20) benchmark accuracy across test suites"])
    ]

    u_w = Inches(5.6)
    u_h = Inches(2.25)

    positions = [
        (Inches(0.8), Inches(1.9)),
        (Inches(6.8), Inches(1.9)),
        (Inches(0.8), Inches(4.5)),
        (Inches(6.8), Inches(4.5))
    ]

    for (title, desc, color, bullets), (ux, uy) in zip(upgrades, positions):
        shape = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ux, uy, u_w, u_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_CARD_BORDER
        shape.line.width = Pt(1.5)

        tb = s4.shapes.add_textbox(ux + Inches(0.2), uy + Inches(0.15), u_w - Inches(0.4), u_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(4)
        p_d.space_after = Pt(4)

        for b in bullets:
            bp = tf.add_paragraph()
            bp.text = b
            bp.font.size = Pt(10)
            bp.font.color.rgb = COLOR_TEXT_MAIN
            bp.space_before = Pt(1)

    try:
        prs.save(output_path)
        print(f"Presentation saved successfully to {output_path}")
    except PermissionError:
        fallback_path = "AI_Finance_Assistant_Demo_DarkMode.pptx"
        prs.save(fallback_path)
        print(f"Original file was open in PowerPoint. Saved dark mode deck to {fallback_path}")

if __name__ == "__main__":
    create_presentation()
