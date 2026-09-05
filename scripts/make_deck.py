"""Generate the demo deck.

    pip install python-pptx
    python scripts/make_deck.py            -> docs/finance-assistant-deck.pptx

Every figure lives in NUMBERS below, so the deck is regenerated from measured
results rather than hand-edited.
"""
from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "finance-assistant-deck.pptx"

# --- measured results; update after re-running the benchmarks -----------------
NUMBERS = {
    "rule_accuracy": "20/20",
    "rule_latency": "35 ms",
    "llm_model": "Qwen2.5-3B",
    "llm_accuracy": "16/20",
    "llm_tokens": "1,960",
    "llm_latency": "18.4 s (CPU)",
    "scale_rows": "20,000,000",
    "scale_size": "3.89 GB",
    "scale_build": "409 s",
    "scale_query": "46 ms",
    "scale_e2e": "401 ms",
    "scale_accuracy": "20/20",
    "guardrail_checks": "14/14",
    "param_cap": "20 B",
    "params_used": "3 B",
}

BG = RGBColor(0x0E, 0x11, 0x17)
PANEL = RGBColor(0x16, 0x1B, 0x25)
LINE = RGBColor(0x26, 0x2E, 0x3D)
TEXT = RGBColor(0xE6, 0xEB, 0xF2)
MUTED = RGBColor(0x8B, 0x96, 0xA8)
ACCENT = RGBColor(0x4A, 0xDE, 0x80)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
FONT = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(s, x, y, w, h, body, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, spacing=1.25):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = body if isinstance(body, list) else [body]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(6)
        content, c, b, sz = line if isinstance(line, tuple) else (line, color, bold, size)
        run = p.add_run()
        run.text = content
        run.font.size = Pt(sz)
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = FONT
    return box


def kicker(s, label):
    text(s, Inches(0.9), Inches(0.55), Inches(11), Inches(0.4), label.upper(), size=12, color=ACCENT, bold=True)


def title(s, heading, sub=None):
    text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(1.0), heading, size=36, bold=True)
    if sub:
        text(s, Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.6), sub, size=17, color=MUTED)


def bullets(s, items, y=Inches(2.7), size=17, gap=0.62, x=Inches(0.9), w=Inches(11.5)):
    for i, item in enumerate(items):
        head, rest = (item if isinstance(item, tuple) else (None, item))
        box = s.shapes.add_textbox(x, y + Inches(i * gap), w, Inches(0.55))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        dot = p.add_run()
        dot.text = "—  "
        dot.font.color.rgb = ACCENT
        dot.font.size = Pt(size)
        dot.font.name = FONT
        if head:
            r = p.add_run()
            r.text = head + "  "
            r.font.bold = True
            r.font.color.rgb = TEXT
            r.font.size = Pt(size)
            r.font.name = FONT
        r = p.add_run()
        r.text = rest
        r.font.color.rgb = MUTED if head else TEXT
        r.font.size = Pt(size)
        r.font.name = FONT


def card(s, x, y, w, h, heading, body, accent=ACCENT):
    box = s.shapes.add_shape(1, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = LINE
    box.line.width = Pt(1)
    box.shadow.inherit = False
    bar = s.shapes.add_shape(1, x, y, Emu(int(Inches(0.05))), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, x + Inches(0.3), y + Inches(0.22), w - Inches(0.5), Inches(0.4), heading, size=15, bold=True)
    text(s, x + Inches(0.3), y + Inches(0.72), w - Inches(0.5), h - Inches(0.9), body, size=12.5, color=MUTED)


def stat(s, x, y, w, value, label, color=ACCENT):
    text(s, x, y, w, Inches(0.8), value, size=40, bold=True, color=color, align=PP_ALIGN.CENTER)
    text(s, x, y + Inches(0.85), w, Inches(0.6), label, size=12.5, color=MUTED, align=PP_ALIGN.CENTER)


def notes(s, body: str) -> None:
    s.notes_slide.notes_text_frame.text = body.strip()


def table(s, rows, x, y, w, col_widths=None, header=True, size=13):
    n_rows, n_cols = len(rows), len(rows[0])
    h = Inches(0.42 * n_rows)
    shape = s.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(0.42)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if (header and r == 0) else BG
            cell.margin_left, cell.margin_right = Inches(0.14), Inches(0.1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if (c and c == n_cols - 1 and not (header and r == 0)) else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.name = FONT
                run.font.bold = header and r == 0
                run.font.color.rgb = MUTED if (header and r == 0) else TEXT
    return shape


def flow(s, steps, y, x=Inches(0.9), w=Inches(11.5), h=Inches(0.85)):
    """A left-to-right chain of boxes with arrows."""
    n = len(steps)
    gap = Inches(0.18)
    bw = Emu(int((w - gap * (n - 1)) / n))
    for i, (label, sub, color) in enumerate(steps):
        bx = x + Emu(int((bw + gap) * i))
        box = s.shapes.add_shape(5, bx, y, bw, h)
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL
        box.line.color.rgb = color
        box.line.width = Pt(1.25)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = FONT
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(10)
        r2.font.color.rgb = MUTED
        r2.font.name = FONT


# ----------------------------------------------------------------- the slides
def build() -> None:
    prs = new_deck()
    n = NUMBERS

    # 1 ---------------------------------------------------------------- title
    s = slide(prs)
    text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.0), "Finance Assistant", size=54, bold=True)
    text(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.8),
         "Ask your ledger a question. Get an answer you can audit.", size=22, color=ACCENT)
    text(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(1.2),
         "The language model reads and writes English.\nSQL does the arithmetic. A verifier checks every digit.",
         size=17, color=MUTED, spacing=1.4)
    text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4),
         f"{n['params_used']} parameter model  ·  {n['param_cap']} cap  ·  {n['scale_rows']} records  ·  every answer traceable to SQL",
         size=13, color=MUTED)
    notes(s, """
Open with the one-liner, do not read the slide.

"Every team here will demo a chatbot that answers finance questions. The difference is what happens
when it does not know. We spent our time on that."

Keep this slide up for 15 seconds. Move.
""")

    # 2 -------------------------------------------------------------- problem
    s = slide(prs)
    kicker(s, "the problem")
    title(s, "Simple questions take hours, not seconds",
          "Finance ops fields the same lookups on repeat, and nobody else can self-serve.")
    bullets(s, [
        ("Dashboards answer the question they were built for.", "Anything else needs a ticket."),
        ("You must know the vocabulary before you can ask.", "Which report? Which field? Which filter?"),
        ("Finance ops becomes a query service desk.", "High-value people running lookups."),
        ("Decisions wait on a queue.", "The answer exists; access to it does not."),
    ])
    notes(s, """
30 seconds. The audience already believes this - do not oversell it.

The line that lands: "the answer already exists in the database. What is missing is access to it."
""")

    # 3 ---------------------------------------------------------------- stakes
    s = slide(prs)
    kicker(s, "why this is different")
    title(s, "A wrong number here is a liability, not a bug",
          "In finance, \"mostly right\" is a failing grade. That single fact drove every design decision.")
    card(s, Inches(0.9), Inches(3.0), Inches(3.6), Inches(2.6), "A plausible number",
         "An LLM will happily produce a total that looks exactly like a real one. Nothing about the output "
         "signals that it was invented.")
    card(s, Inches(4.75), Inches(3.0), Inches(3.6), Inches(2.6), "Corrupts what it touches",
         "It flows into reconciliation, audit trails and board reporting before anyone checks it.", AMBER)
    card(s, Inches(8.6), Inches(3.0), Inches(3.6), Inches(2.6), "Destroys trust permanently",
         "One fabricated figure and the assistant is never used again. Correctness is the product.",
         RGBColor(0xF8, 0x71, 0x71))
    notes(s, """
This is the slide that frames everything after it. Slow down.

"A chatbot that is 95% right is a useful product. A finance assistant that is 95% right is a
liability, because you cannot tell which 5% is wrong - and it looks exactly like the other 95%."

Then: "so we did not build a chatbot with guardrails bolted on. We built a calculator that happens
to speak English."
""")

    # 4 --------------------------------------------------------------- thesis
    s = slide(prs)
    kicker(s, "our approach")
    title(s, "Ask the model to do only what models are good at",
          "It never sees a number, never writes SQL, never does arithmetic.")
    flow(s, [
        ("Question", "plain language", BLUE),
        ("Query plan", "3B model, JSON only", BLUE),
        ("Validate", "whitelist schema", ACCENT),
        ("SQL", "built in Python", ACCENT),
        ("DuckDB", "computes the number", ACCENT),
        ("Verify", "every digit checked", AMBER),
    ], y=Inches(3.0))
    text(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(1.6),
         [("The model's only output is a fixed 10-key JSON plan: dataset, filters, grouping, period.", TEXT, True, 17),
          ("There is no field in that schema where a figure can be written. Hallucinating a number is not "
           "something we detect and correct — it is something the interface makes impossible.", MUTED, False, 15)],
         spacing=1.35)
    notes(s, """
The key sentence is the bold one. Say it, then pause.

"The model fills in a form. The form has ten fields: dataset, filters, grouping, period. None of
them is a number. So when people ask how we stop it hallucinating a total - we did not stop it.
We removed the place where it could write one."

If asked why not text-to-SQL: free-form SQL from a small model is unverifiable. You cannot prove
the query means what the user asked. A fixed plan schema is checkable field by field - and that is
precisely what makes a 3B model sufficient.
""")

    # 5 --------------------------------------------------------- architecture
    s = slide(prs)
    kicker(s, "architecture")
    title(s, "Three tiers of understanding, one path to a number")
    rows = [
        ["Stage", "Who does it", "Why"],
        ["Understand the question", "LLM (3B)", "Language is what models are good at"],
        ["Choose filters, grouping, period", "LLM → JSON schema", "Tiny output space, fully validatable"],
        ["Resolve dates", "Python", "\"Last month\" must be exact, not probable"],
        ["Resolve counterparty names", "Fuzzy match vs real data", "The model cannot invent a payee"],
        ["Build SQL", "Whitelist builder", "Identifiers from catalog, literals bound"],
        ["Filter, group, aggregate", "DuckDB", "Arithmetic is exact and reproducible"],
        ["Repeat aggregates", "Pre-built monthly rollup", "Summary questions skip the raw scan"],
        ["Phrase the answer", "LLM (3B)", "Reads only pre-computed facts"],
        ["Publish", "Number guardrail", "Untraceable digit → answer replaced"],
    ]
    table(s, rows, Inches(0.9), Inches(2.4), Inches(11.5),
          col_widths=[Inches(3.6), Inches(3.2), Inches(4.7)], size=12)
    notes(s, """
Do not read the table. Point at the two green rows in the middle and say:

"Everything from here down is deterministic Python and SQL. The model is involved at the top and
at the bottom - understanding the question, and phrasing the result. Never in between."

Worth noting: the whitelist builder means SQL injection is not mitigated, it is structurally
impossible. Identifiers come from a catalog, every literal is a bound parameter.
""")

    # 6 ------------------------------------------------------------ guardrails
    s = slide(prs)
    kicker(s, "grounding")
    title(s, "Five guardrails. Each one alone would stop a fabricated figure.")
    items = [
        ("Cannot express a number.", "Its only output is a query plan. No field holds a figure."),
        ("Cannot express a column.", "Invented fields are stripped by the validator before SQL is built."),
        ("Cannot express an entity.", "Names are matched against the data. No match → it says so."),
        ("Cannot compute.", "Dates resolved in Python, maths in DuckDB, connection read-only."),
        ("Its wording is verified.", "Every digit traced to the result set. One miss → sentence discarded."),
    ]
    bullets(s, items, y=Inches(2.6), gap=0.6, size=16)
    card(s, Inches(0.9), Inches(5.85), Inches(11.5), Inches(1.0), "And when there is no answer",
         "Unknown counterparty, unsupported question or zero matching rows produce an explicit refusal — "
         "never an estimate. Two of the twenty benchmark questions exist purely to test this.", AMBER)
    notes(s, """
Go fast on the five - one line each, they are self-explanatory.

Spend the time on the amber box instead. "Two of our twenty benchmark questions are designed to be
refused. Most teams will not have tested refusal at all, because it does not demo well. It is the
only part of this we would stake our reputation on."
""")

    # 7 ------------------------------------------------------------- the data
    s = slide(prs)
    kicker(s, "the data")
    title(s, "The schema has no counterparty and no reconciliation column",
          "Both are central to the brief. We derived them from real columns rather than inventing them.")
    rows = [
        ["Field", "Derived from", "Stated to the user"],
        ["counterparty", "Longest capitalised run in the narration, rail noise stripped",
         "\"parsed out of free-text\""],
        ["channel", "UPI / NEFT / IMPS / RTGS / FT / CHEQUE / ACH / CHARGES", "shown as a field"],
        ["reconciliation_status", "Has a bank reference OR a UTR → reconciled", "definition quoted in full"],
        ["account_number_masked", "Last 4 digits only", "raw column not in the catalog"],
    ]
    table(s, rows, Inches(0.9), Inches(2.9), Inches(11.5),
          col_widths=[Inches(2.7), Inches(5.6), Inches(3.2)], size=12)
    card(s, Inches(0.9), Inches(5.2), Inches(5.6), Inches(1.6), "Validated against real narrations",
         "All 10 production sample rows in the schema doc parse to the right counterparty and channel, "
         "and the one row with neither reference nor UTR is correctly flagged unreconciled. Asserted on every test run.")
    card(s, Inches(6.8), Inches(5.2), Inches(5.6), Inches(1.6), "Reads your real database",
         "A sync job pulls bank / account / transaction straight from MySQL, decrypts the protected columns "
         "on the way in, and materialises the analysis copy. No manual export step.", BLUE)
    notes(s, """
This slide exists because a sharp judge will ask "where did 'counterparty' come from? That is not
in the schema you were given."

Answer: "Correct. Neither is reconciliation. We had two options - drop half the brief, or derive
them. We derived them from real columns, wrote the definition down once, and the assistant states
that definition in every answer that depends on it. It is derived, documented and visible - not
assumed."

On the reference-vs-UTR question the schema doc raises: a bare "reference number" hits
transaction_reference_id, the plaintext searchable column. UTR only decides reconciliation state.
""")

    # 8 ---------------------------------------------------------- security
    s = slide(prs)
    kicker(s, "sensitive data")
    title(s, "The schema flags two columns as sensitive. Neither can leave the database.",
          "account_number and utr_number are protected at four independent layers.")
    rows = [
        ["Layer", "What it does"],
        ["Encrypted at rest", "AES-256-SIV. Deterministic, so encrypted columns still filter and group."],
        ["Decrypted once, at load", "Never at query time. Persisted views are pure SQL, so the connection stays read-only."],
        ["Not in the query catalog", "account_number is exposed only masked; utr_number is not exposed at all."],
        ["Refused at the planner", "\"Show me the full UTR\" returns a security refusal, not a query."],
    ]
    table(s, rows, Inches(0.9), Inches(2.9), Inches(11.5),
          col_widths=[Inches(3.4), Inches(8.1)], size=13)
    card(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.4), "Why four layers and not one",
         "The planner refusal is the weakest of them, because it depends on the model behaving. The other "
         "three do not: even a model that explicitly asks for utr_number cannot get it, because the column "
         "is not in the catalog the SQL builder is allowed to name. There is a test that proves exactly this.",
         AMBER)
    notes(s, """
The schema doc calls out account_number and utr_number as sensitive and asks teams to decide how to
handle them. This slide is our answer.

The honest point to make is the amber box: "the planner refusal is the layer we trust least, because
it depends on the model doing what it is told. It is there for good manners. The real protection is
that the column does not exist as far as the SQL builder is concerned."

On the reference-vs-UTR question the doc raises: a bare "reference number" resolves to
transaction_reference_id, the plaintext searchable column. UTR is used only to derive reconciliation
state. That decision is made once, in code, not per question.
""")

    # 9 ---------------------------------------------------- model choice
    s = slide(prs)
    kicker(s, "model choice")
    title(s, f"{n['llm_model']} against a {n['param_cap']} cap — and a 0-parameter tier that already passes",
          "We got the accuracy from architecture, not from scale.")
    rows = [
        ["Tier", "What runs", "When", "Accuracy"],
        ["1", f"{n['llm_model']}, JSON mode, temp 0", "Every question", n["llm_accuracy"]],
        ["2", "Same model, one repair round", "Plan fails validation", "—"],
        ["3", "Deterministic rule parser", "Model unreachable or invalid", n["rule_accuracy"]],
    ]
    table(s, rows, Inches(0.9), Inches(2.9), Inches(11.5),
          col_widths=[Inches(0.9), Inches(4.2), Inches(3.9), Inches(2.5)])
    text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.4),
         [("Tier 3 is not a stub — it scores full marks on its own.", TEXT, True, 16),
          ("That is the honest way to show how little of the accuracy comes from model scale. The LLM widens "
           "the range of phrasings we understand; it does not produce the numbers. If it goes down mid-demo, "
           "every question still answers correctly.", MUTED, False, 14)],
         spacing=1.3)
    notes(s, """
This is the slide the "lowest possible model" criterion is scored on. Do not rush it.

"The brief says smallest model, highest accuracy. Most teams will answer that by picking a small
model and hoping. We answered it by measuring what the model actually contributes - and the answer
is that our zero-parameter tier already passes the benchmark."

"So the 3B is not carrying the accuracy. It is carrying the language understanding. That is the
only job we could not do deterministically, and it is a job a 3B model is genuinely good at."

If the LLM tier number is below the rule tier, say so plainly and explain: the LLM tier handles
phrasings the rules do not cover; the rule tier is the floor, not the ceiling.
""")

    # 10 --------------------------------------------------------------- numbers
    s = slide(prs)
    kicker(s, "measured, not claimed")
    title(s, "Accuracy and scale")
    stat(s, Inches(0.9), Inches(2.5), Inches(2.6), n["rule_accuracy"], "benchmark vs ground-truth SQL")
    stat(s, Inches(3.7), Inches(2.5), Inches(2.6), n["guardrail_checks"], "anti-hallucination checks", BLUE)
    stat(s, Inches(6.5), Inches(2.5), Inches(2.6), n["scale_e2e"], f"end to end at {n['scale_rows']} rows", BLUE)
    stat(s, Inches(9.3), Inches(2.5), Inches(2.6), n["params_used"], f"parameters (cap {n['param_cap']})", AMBER)
    rows = [
        ["At the 20 M-record test limit", ""],
        ["Database file", n["scale_size"]],
        ["One-off build: generate, derive, index", n["scale_build"]],
        ["Single-month aggregate", n["scale_query"]],
        ["End-to-end, question to answer", n["scale_e2e"]],
        ["Accuracy at 20 M vs 250 k", f"{n['scale_accuracy']} — unchanged"],
    ]
    table(s, rows, Inches(0.9), Inches(4.4), Inches(6.5), col_widths=[Inches(4.6), Inches(1.9)])
    card(s, Inches(7.9), Inches(4.4), Inches(4.5), Inches(2.5), "Derived, not invented",
         "The schema has no counterparty and no reconciliation column. Both are derived from real columns "
         "with one documented definition — and the assistant states that definition in every answer that "
         "uses it. All 10 production sample narrations parse correctly.", BLUE)
    notes(s, """
Lead with: "these are measured on a laptop this morning, not estimated."

The number to emphasise is the last row of the table - accuracy is unchanged between 250 thousand
and 20 million rows. Scale did not cost us correctness.

If asked how: the queries are vectorised scans over two or three columns with a date predicate.
The expensive part - parsing narrations - happens once at load, never per query. And we moved the
top-N record sample off the answer path because profiling showed it was 80% of the latency.
""")

    # 11 ------------------------------------------------------------ demo flow
    s = slide(prs)
    kicker(s, "demo")
    title(s, "Five moments, four minutes")
    demo = [
        ("1.  \"How much did we pay out last month?\"",
         "Then open \"How I got this answer\": the plan the model wrote, the SQL we built, the source records."),
        ("2.  \"How does that compare to the month before?\"",
         "No context repeated. Multi-turn carries the previous plan forward."),
        ("3.  \"Which transactions are still unreconciled?\"",
         "Download to Excel — the finance user leaves with the working file, not a screenshot."),
        ("4.  \"...with Globex Corporation?\"  /  \"...next quarter?\"",
         "Both refused, out loud. The 60 seconds we most want you to watch."),
        ("5.  scripts/selftest.py",
         "A stubbed model claims a figure that is nowhere in the data. Watch the verifier reject it."),
    ]
    bullets(s, demo, y=Inches(2.6), gap=0.72, size=15)
    text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
         "If the model host dies mid-demo, nothing breaks — the rule tier answers all of these correctly.",
         size=13, color=AMBER)
    notes(s, """
RUN THIS ON THE LAPTOP WHERE OLLAMA WORKS. Check before you start:
  - ollama is serving, and the sidebar says "LLM planner + narrator" not "rule parser"
  - data/finance.duckdb was rebuilt with the CURRENT code (the rollup table must exist)
  - port 8000 is free

Switch to the browser now. Do not talk over the demo.

Moment 1: point at the confidence badge, the record count and the token counter. Then open the
explain panel and scroll it slowly.
Moment 4 is the one that wins this. Say out loud: "most demos avoid these two questions. This is
the one we most want you to ask."
Moment 5: drop to the terminal. The stubbed model claims 42 million. Watch it get rejected and the
confidence drop to medium.

If anything breaks: the sidebar shows whether the LLM is reachable. If it is not, say so and keep
going - it still works. That is the strongest possible statement about where the accuracy lives.
""")

    # 12 --------------------------------------------------------------- impact
    s = slide(prs)
    kicker(s, "impact")
    title(s, "Seconds instead of a ticket — and an audit trail either way")
    card(s, Inches(0.9), Inches(2.7), Inches(3.6), Inches(2.4), "For the business",
         "Anyone can ask. Finance ops stops being a query service desk and goes back to higher-value work.")
    card(s, Inches(4.75), Inches(2.7), Inches(3.6), Inches(2.4), "For audit",
         "Every answer ships with the plan, the SQL, the assumptions and the source rows. Reproducible on demand.",
         BLUE)
    card(s, Inches(8.6), Inches(2.7), Inches(3.6), Inches(2.4), "For security",
         "Sensitive columns encrypted at rest and unreachable by query. A 3B model on a laptop — no "
         "frontier API bill, and no financial data leaving the building.", AMBER)
    text(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.2),
         [("Next:", TEXT, True, 16),
          ("more datasets behind the same plan schema  ·  scheduled digests  ·  approval workflows  ·  "
           "the same guardrail in front of any model", MUTED, False, 15)], spacing=1.3)
    notes(s, """
Close on the audit line, not the cost line.

"The reason this matters is not that it saves a finance analyst ten minutes. It is that every
answer it gives can be reproduced from the plan and the SQL it shows you. That is the difference
between a demo and something you would actually let near a reconciliation."

Then stop talking.
""")

    # 13 ------------------------------------------------------------------ q&a
    s = slide(prs)
    kicker(s, "backup")
    title(s, "Questions we expect")
    qa = [
        ("\"What if the model picks the wrong filter?\"",
         "Visible in the explain panel, and asserted structurally in the benchmark - not just the final number."),
        ("\"Why not text-to-SQL?\"",
         "Free-form SQL from a small model is unverifiable. A fixed plan schema is checkable field by field."),
        ("\"Does it scale?\"",
         f"{n['scale_rows']} rows, {n['scale_e2e']} end to end, accuracy unchanged. Measured, reproducible."),
        ("\"What about a question you don't support?\"",
         "It refuses and says what it can answer. Coverage grows by adding fields to one catalog file."),
        ("\"Is the LLM doing anything at all?\"",
         "Language understanding. Turn it off and the scripted questions still pass - that is the point."),
        ("\"How do you handle the sensitive columns?\"",
         "Encrypted at rest, decrypted once at load, absent from the query catalog, refused at the planner."),
    ]
    bullets(s, qa, y=Inches(2.4), gap=0.72, size=14.5)
    notes(s, """
Backup slide - only show it if Q&A goes long or you are asked something on this list.
""")

    OUT.parent.mkdir(exist_ok=True)
    target = OUT
    try:
        prs.save(target)
    except PermissionError:  # the deck is open in PowerPoint
        target = OUT.with_name(f"{OUT.stem}-{datetime.now():%H%M%S}{OUT.suffix}")
        prs.save(target)
        print(f"{OUT.name} is open elsewhere - wrote a copy instead")
    print(f"{len(prs.slides._sldIdLst)} slides -> {target}")


if __name__ == "__main__":
    build()
